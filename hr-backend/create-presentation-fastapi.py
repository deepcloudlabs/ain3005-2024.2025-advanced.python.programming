from pptx import Presentation

# Create a PowerPoint presentation
presentation = Presentation()

# Slide titles and contents for the FastAPI presentation
slides_content = [
    ("Introduction to FastAPI", "FastAPI is a modern, fast web framework for building APIs with Python 3.6+.\nKey Features:\n- High performance\n- Easy to use\n- Based on standard Python type hints."),
    ("Why Choose FastAPI?", "Benefits:\n- Fast to code\n- High performance\n- Automatic interactive API documentation\n- Based on OpenAPI and JSON Schema."),
    ("Core Features", "Key Features:\n- Request validation\n- Dependency injection\n- Asynchronous programming support\n- Automatic documentation with Swagger and ReDoc."),
    ("Installation", "Install FastAPI and Uvicorn (ASGI server):\n\n```bash\npip install fastapi uvicorn\n```"),
    ("Creating Your First API", "1. Import FastAPI:\n```python\nfrom fastapi import FastAPI\n```\n2. Create an instance:\n```python\napp = FastAPI()\n```\n3. Define routes:\n```python\n@app.get('/')\ndef read_root():\n    return {'Hello': 'World'}\n```"),
    ("Running the API", "Run the app using Uvicorn:\n\n```bash\nuvicorn main:app --reload\n```"),
    ("Path Parameters", "Define parameters in the URL path:\n\n```python\n@app.get('/items/{item_id}')\ndef read_item(item_id: int):\n    return {'item_id': item_id}\n```"),
    ("Query Parameters", "Add optional query parameters:\n\n```python\n@app.get('/items/')\ndef read_items(skip: int = 0, limit: int = 10):\n    return {'skip': skip, 'limit': limit}\n```"),
    ("Request Body", "Define the request body using Pydantic models:\n\n```python\nfrom pydantic import BaseModel\nclass Item(BaseModel):\n    name: str\n    price: float\n\n@app.post('/items/')\ndef create_item(item: Item):\n    return item\n```"),
    ("Dependency Injection", "Use dependencies to share logic across routes:\n\n```python\nfrom fastapi import Depends\n\ndef common_params(q: str = None):\n    return q\n\n@app.get('/items/')\ndef read_items(q: str = Depends(common_params)):\n    return {'q': q}\n```"),
    ("Async Programming", "Leverage Python's `async` capabilities:\n\n```python\n@app.get('/async/')\nasync def read_async():\n    return {'message': 'This is async!'}\n```"),
    ("Error Handling", "Use `HTTPException` to handle errors:\n\n```python\nfrom fastapi import HTTPException\n\n@app.get('/items/{item_id}')\ndef read_item(item_id: int):\n    if item_id > 10:\n        raise HTTPException(status_code=404, detail='Item not found')\n    return {'item_id': item_id}\n```"),
    ("Middleware", "Add middleware to intercept requests and responses:\n\n```python\nfrom fastapi.middleware.cors import CORSMiddleware\n\napp.add_middleware(\n    CORSMiddleware,\n    allow_origins=['*'],\n    allow_credentials=True,\n    allow_methods=['*'],\n    allow_headers=['*']\n)\n```"),
    ("Static Files", "Serve static files using `StaticFiles`:\n\n```python\nfrom fastapi.staticfiles import StaticFiles\n\napp.mount('/static', StaticFiles(directory='static'), name='static')\n```"),
    ("Authentication", "FastAPI supports OAuth2 and JWT:\n- Use `fastapi.security` for authentication schemes.\n- Example:\n\n```python\nfrom fastapi.security import OAuth2PasswordBearer\n\noauth2_scheme = OAuth2PasswordBearer(tokenUrl='token')\n@app.get('/users/me/')\ndef read_users_me(token: str = Depends(oauth2_scheme)):\n    return {'token': token}\n```"),
    ("Swagger UI", "Interactive API docs available by default:\n- Swagger UI: [http://127.0.0.1:8000/docs](http://127.0.0.1:8000/docs)\n- ReDoc: [http://127.0.0.1:8000/redoc](http://127.0.0.1:8000/redoc)"),
    ("Testing", "FastAPI supports testing with `TestClient`:\n\n```python\nfrom fastapi.testclient import TestClient\n\nclient = TestClient(app)\n\ndef test_read_main():\n    response = client.get('/')\n    assert response.status_code == 200\n    assert response.json() == {'Hello': 'World'}\n```"),
    ("Performance", "FastAPI is built on Starlette and Pydantic:\n- Starlette: High-performance ASGI framework.\n- Pydantic: Data validation and settings management."),
    ("Community and Ecosystem", "FastAPI has a growing ecosystem:\n- Third-party libraries\n- Integrations with databases and message queues\n- Active community and documentation."),
    ("Conclusion", "FastAPI is a robust framework for building modern APIs:\n- Focuses on performance and developer experience.\n- Easy to learn, scale, and integrate.\nTry it today!")
]

# Add slides to the presentation
for title, content in slides_content:
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])  # Title and Content slide layout
    slide.shapes.title.text = title
    slide.placeholders[1].text = content

# Save the presentation
file_path = "FastAPI_Tutorial.pptx"
presentation.save(file_path)